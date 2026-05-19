from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Salesforce brand colors
SF_BLUE = RGBColor(0x00, 0x6D, 0xFF)   # #006DFF
SF_DARK = RGBColor(0x03, 0x2D, 0x60)   # #032D60
SF_TEAL = RGBColor(0x06, 0xA5, 0x9A)   # #06A59A
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY   = RGBColor(0xD0, 0xD7, 0xE5)
DARK_GRAY  = RGBColor(0x3E, 0x3E, 0x3C)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=11, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Salesforce Sans" if False else "Calibri"
    return txBox

def add_bullet_box(slide, items, left, top, width, height,
                   font_size=10, color=DARK_GRAY, bullet_char="▸ "):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

# ---- Build slide ----
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

W = 13.33
H = 7.5

# --- Background ---
add_rect(slide, 0, 0, W, H, fill_color=LIGHT_GRAY)

# --- Header bar ---
add_rect(slide, 0, 0, W, 1.1, fill_color=SF_DARK)

# Title
add_text_box(slide, "Webhook Bridge  |  Tableau Extension", 0.35, 0.12, 8, 0.55,
             font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle tagline
add_text_box(slide, "Turn any Tableau dashboard into a trigger for operational workflows",
             0.35, 0.63, 9, 0.38,
             font_size=12, bold=False, color=rgb(0xB0, 0xC9, 0xFF), align=PP_ALIGN.LEFT)

# Logo area placeholder (top right)
add_text_box(slide, "Salesforce SE", 11.0, 0.2, 2.1, 0.4,
             font_size=10, bold=False, color=rgb(0x80, 0xA8, 0xE0), align=PP_ALIGN.RIGHT)

# --- Three column cards ---
card_top    = 1.28
card_height = 3.6
card_w      = 3.85
gap         = 0.19
col1 = 0.28
col2 = col1 + card_w + gap
col3 = col2 + card_w + gap

for cx in [col1, col2, col3]:
    add_rect(slide, cx, card_top, card_w, card_height,
             fill_color=WHITE, line_color=MID_GRAY, line_width=Pt(0.5))

# ---- Card 1: The Problem ----
add_rect(slide, col1, card_top, card_w, 0.38, fill_color=SF_DARK)
add_text_box(slide, "THE GAP", col1 + 0.12, card_top + 0.05, card_w - 0.2, 0.3,
             font_size=10, bold=True, color=WHITE)

add_text_box(slide, "Tableau shows the insight. Then what?",
             col1 + 0.15, card_top + 0.52, card_w - 0.3, 0.45,
             font_size=11.5, bold=True, color=SF_DARK)

add_bullet_box(slide, [
    "Users copy data out manually to act on it",
    "No native way to trigger external workflows from a mark selection",
    "Tableau Flow actions only work for Salesforce — not Zapier, Make, Power Automate, n8n, MuleSoft, or custom endpoints",
    "Every action requires leaving the dashboard",
], col1 + 0.15, card_top + 1.05, card_w - 0.3, 2.4,
   font_size=10, color=DARK_GRAY)

# ---- Card 2: The Solution ----
add_rect(slide, col2, card_top, card_w, 0.38, fill_color=SF_TEAL)
add_text_box(slide, "THE SOLUTION", col2 + 0.12, card_top + 0.05, card_w - 0.2, 0.3,
             font_size=10, bold=True, color=WHITE)

add_text_box(slide, "Select. Click. Done.",
             col2 + 0.15, card_top + 0.52, card_w - 0.3, 0.45,
             font_size=11.5, bold=True, color=SF_DARK)

add_bullet_box(slide, [
    "Dashboard authors configure a webhook URL + field mappings — once",
    "Viewers select marks and click a button; data POSTs as structured JSON",
    "Works with any HTTP endpoint: Zapier, Make, Power Automate, n8n, MuleSoft",
    "No backend, no code — deploys in 10 minutes via GitHub Pages",
    "Config travels with the workbook",
], col2 + 0.15, card_top + 1.05, card_w - 0.3, 2.4,
   font_size=10, color=DARK_GRAY)

# ---- Card 3: The Opportunity ----
add_rect(slide, col3, card_top, card_w, 0.38, fill_color=SF_BLUE)
add_text_box(slide, "THE OPPORTUNITY", col3 + 0.12, card_top + 0.05, card_w - 0.2, 0.3,
             font_size=10, bold=True, color=WHITE)

add_text_box(slide, "Every Tableau customer. Any automation stack.",
             col3 + 0.15, card_top + 0.52, card_w - 0.3, 0.45,
             font_size=11.5, bold=True, color=SF_DARK)

add_bullet_box(slide, [
    "Positions Tableau as a control surface for operations, not just analysis",
    "Relevant to any account with an automation stack (virtually all enterprise)",
    "Low barrier: free, no install, no admin approval required",
    "Demo-ready: live webhook can be shown in 10 minutes with webhook.site",
    "Opens conversation about Tableau + MuleSoft / Flow / Agentforce integrations",
], col3 + 0.15, card_top + 1.05, card_w - 0.3, 2.4,
   font_size=10, color=DARK_GRAY)

# --- How it works row ---
how_top = card_top + card_height + 0.18
add_rect(slide, 0.28, how_top, W - 0.56, 0.32, fill_color=MID_GRAY)
add_text_box(slide, "HOW IT WORKS", 0.42, how_top + 0.05, 4, 0.25,
             font_size=9, bold=True, color=SF_DARK)

steps = [
    ("1  Add Extension", "Drag extension zone onto dashboard, load .trex manifest"),
    ("2  Configure",      "Set webhook URL, worksheet, field mappings — saved to workbook"),
    ("3  Select Marks",   "Dashboard viewer selects data points of interest"),
    ("4  Click Button",   "Structured JSON POSTs to Zapier / Make / Power Automate / etc."),
    ("5  Workflow Runs",  "Downstream automation does whatever you need — CRM update, alert, ticket..."),
]

step_w = (W - 0.56) / len(steps)
for i, (title, desc) in enumerate(steps):
    sx = 0.28 + i * step_w
    sy = how_top + 0.32
    fill = SF_DARK if i % 2 == 0 else rgb(0x06, 0x5A, 0xA5)
    add_rect(slide, sx, sy, step_w - 0.02, 1.25, fill_color=fill)
    add_text_box(slide, title, sx + 0.1, sy + 0.1, step_w - 0.2, 0.3,
                 font_size=9.5, bold=True, color=WHITE)
    add_text_box(slide, desc, sx + 0.1, sy + 0.38, step_w - 0.2, 0.8,
                 font_size=8.5, bold=False, color=rgb(0xB8, 0xD4, 0xFF))

# --- Footer ---
footer_top = H - 0.28
add_rect(slide, 0, footer_top, W, 0.28, fill_color=SF_DARK)
add_text_box(slide, "Working prototype · Tested on Tableau Cloud · github.com/bhartSF/tableau-webhook-bridge",
             0.3, footer_top + 0.05, W - 0.6, 0.22,
             font_size=8, bold=False, color=rgb(0x80, 0xA8, 0xE0), align=PP_ALIGN.LEFT)
add_text_box(slide, "Ben Hart · Salesforce SE",
             W - 3.0, footer_top + 0.05, 2.8, 0.22,
             font_size=8, bold=False, color=rgb(0x80, 0xA8, 0xE0), align=PP_ALIGN.RIGHT)

out_path = "/Users/ben.hart/Documents/Coding Projects/claude-projects/tableau-native-action-framework/webhook_bridge_pitch.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
